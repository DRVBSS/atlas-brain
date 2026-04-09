"""PowerPoint extraction via python-pptx."""

from pathlib import Path

from atlas_brain.models import ProcessedDocument


def extract(file_path: Path) -> ProcessedDocument:
    """Extract text from PPTX files (slide text + speaker notes)."""
    from pptx import Presentation

    prs = Presentation(str(file_path))

    slides_text = []
    sections = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"\n*Speaker notes:* {notes}")

        slide_text = "\n".join(slide_parts)
        slides_text.append(slide_text)
        sections.append(f"Slide {i}")

    full_text = "\n\n".join(slides_text)
    title = file_path.stem

    # Try first slide title
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title = text[:100]
                    break

    return ProcessedDocument(
        text=full_text,
        title=title,
        author=None,
        created_date=None,
        word_count=len(full_text.split()),
        sections=sections,
        metadata={"slide_count": len(prs.slides)},
    )
